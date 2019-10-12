from pathlib import Path
from typing import List

from pptx import Presentation
from pptx.shapes.autoshape import BaseShape
from pptx.slide import Slide

example_file_path = Path("./data/example01.pptx").resolve()
prs = Presentation(example_file_path)
print(prs)


def get_shape_info(shape: BaseShape) -> dict:
    return {
        "name": shape.name,
        "text": shape.text if shape.has_text_frame else None,
        "width": shape.width,
        "height": shape.height,
        "id": shape.shape_id,
        "x": shape.left,
        "y": shape.top,
        "type": type(shape),
    }


slide_info = [get_shape_info(shape) for shape in prs.slides[0].shapes]
print(slide_info)


def replace_by_image(slide: Slide, name: str, img_path: Path, *, do_not_scale: bool = False) -> List[BaseShape]:
    shapes_to_replace = [shape for shape in slide.shapes if hasattr(shape, "text") and shape.text == name]
    print(len(shapes_to_replace))
    new_image_shapes = []
    for old_shape in shapes_to_replace:
        shape_info = get_shape_info(old_shape)
        print(shape_info)
        img_file = open(img_path, "rb")
        slide_shapes = old_shape._parent
        img_shape = slide_shapes.add_picture(
            img_file,
            old_shape.left,
            old_shape.top,
        )
        old_aspect_ratio = old_shape.width / old_shape.height
        new_aspect_ratio = img_shape.width / img_shape.height
        if img_shape.height <= old_shape.height and img_shape.width <= old_shape.width and not do_not_scale:
            if old_aspect_ratio >= new_aspect_ratio:
                img_shape.width = old_shape.width
                img_shape.height = int(img_shape.width / new_aspect_ratio)
            else:
                img_shape.height = old_shape.height
                img_shape.width = int(img_shape.height * new_aspect_ratio)
        img_shape.top += int((old_shape.height - img_shape.height) / 2)
        img_shape.left += int((old_shape.width - img_shape.width) / 2)
        new_image_shapes.append(img_shape)
        slide_shapes.element.remove(old_shape.element)
    return new_image_shapes


added_img_shapes = replace_by_image(prs.slides[0], "#logo", Path("./playground/pptx_icon.png"))
assert len(added_img_shapes) == 1
added_img_shapes = replace_by_image(prs.slides[0], "#logo", Path("./playground/pptx_icon.png"))
assert len(added_img_shapes) == 0
prs.save(Path("./playground") / example_file_path.name)
