import pathlib
import pptx
from typing import Union, Iterable, Tuple
from pptx.shapes.base import BaseShape

_Pathlike = Union[str, pathlib.Path]


class Template:
    """Helper class for modifying pptx templates.
    """

    def __init__(self, filename: _Pathlike) -> None:
        """Initializes a Template-Modifier.

        Args:
            filename (path-like): file name or path
        """
        self._template_path = filename
        self._presentation = pptx.Presentation(filename)
        pass

    def replace_text(self, label: str, new_text: str) -> None:
        """Replaces text placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            text (str): new content
            scope: None, slide number, Slide object or iterable of Slide objects
        """
        slide_number, tag_name = self._parse_label(label)
        shapes = self._find_shapes(slide_number, tag_name)
        for shape in shapes:
            shape.text = new_text

    def replace_picture(self, label: str, filename: _Pathlike, *, scale_up: bool = True) -> None:
        """Replaces rectangle placeholders on one or many slides.

        The aspect ratio of the image is not changed.
        To large images are always resized.
        The behaviour for small images is configurable.
        Centering is always active.

        Args:
            label (str): label of the placeholder (without curly braces)
            filename (path-like): path to an image file
            scale_up (bool): deactivates that the image is enlarged (default: True)
        """
        slide_number, tag_name = self._parse_label(label)
        shapes_to_replace = self._find_shapes(slide_number, tag_name)
        if not shapes_to_replace:
            raise ValueError(f"The label '{label}' can't be found in the template.")
        with open(filename, "rb") as img_file:
            old_shape: BaseShape
            for old_shape in shapes_to_replace:
                slide_shapes = old_shape._parent
                img_shape = slide_shapes.add_picture(
                    image_file=img_file,
                    left=old_shape.left,
                    top=old_shape.top,
                )
                # Scaling the image if `scale_up == True`:
                if img_shape.height <= old_shape.height and img_shape.width <= old_shape.width and scale_up:
                    old_aspect_ratio = old_shape.width / old_shape.height
                    new_aspect_ratio = img_shape.width / img_shape.height
                    if old_aspect_ratio >= new_aspect_ratio:
                        img_shape.width = old_shape.width
                        img_shape.height = int(img_shape.width / new_aspect_ratio)
                    else:
                        img_shape.height = old_shape.height
                        img_shape.width = int(img_shape.height * new_aspect_ratio)
                # Centering the image at the extent of the placeholder:
                img_shape.top += int((old_shape.height - img_shape.height) / 2)
                img_shape.left += int((old_shape.width - img_shape.width) / 2)
                del slide_shapes[slide_shapes.index(old_shape)]
                # Removing shapes is performed at the lxml level.
                # The `element` attribute contains an instance of `lxml.etree._Element`.
                slide_shapes.element.remove(old_shape.element)

    def replace_table(self, label: str, data) -> None:
        """Replaces rectangle placeholders on one or many slides.

        Args:
            label (str): label of the placeholder (without curly braces)
            data (pandas.DataFrame): table to be inserted into the presentation
        """
        pass

    def _parse_label(self, label: str) -> Tuple[Union[int, str], str]:
        slide_number, tag_name = label.split(':')
        return int(slide_number) if slide_number != '*' else slide_number, tag_name

    def _find_shapes(self,
                     slide_number: Union[int, str],
                     tag_name: str) -> Iterable[BaseShape]:
        """Finds all shapes that match the label

        Args:
            label (str): label of the placeholder (without curly braces)
        """
        matched_shapes = []

        def _find_shapes_in_slide(slide):
            return filter(lambda shape: shape.text == f'{{{tag_name}}}', slide.shapes)

        if slide_number == '*':
            slides = self._presentation.slides
        else:
            # in label we are using 1 based indexing
            slide_index = slide_number - 1
            if slide_index < 0 or slide_index >= len(self._presentation.slides):
                raise IndexError(f"Can't find slide number {slide_number}.")
            slides = [self._presentation.slides[slide_index]]

        for slide in slides:
            matched_shapes.extend(_find_shapes_in_slide(slide))

        return matched_shapes

    def save(self, filename: _Pathlike) -> None:
        """Saves the updated pptx to the specified filepath.

        Args:
            filename (path-like): file name or path
        """
        # TODO: make sure that the user does not override the self._template_path
        pass
